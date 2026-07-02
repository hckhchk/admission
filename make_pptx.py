from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

OUT = r'C:\Users\user\Desktop\yes\입시대시보드_연수자료.pptx'

# ── 팔레트 ─────────────────────────────────
def rgb(h): return RGBColor.from_string(h)

NAVY   = rgb("1A2A4A"); BLUE    = rgb("2E5C9E"); LT_BLUE = rgb("CADCFC")
GOLD   = rgb("D4950A"); WHITE   = rgb("FFFFFF"); LT_BG   = rgb("F0F4FA")
DK_TXT = rgb("1A1A2E"); GRAY    = rgb("64748B"); LT_GRAY = rgb("D1D9E6")
RED    = rgb("C0392B"); ORANGE  = rgb("E67E22"); GREEN   = rgb("1A7A4A")

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

# ── 헬퍼 ───────────────────────────────────
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def rect(s, x, y, w, h, fc, lc=None, lw=Pt(0)):
    sh = s.shapes.add_shape(MSO.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if lc: sh.line.color.rgb = lc; sh.line.width = lw
    else:  sh.line.fill.background()
    return sh

def rrect(s, x, y, w, h, fc, lc=None, lw=Pt(0.75), adj=0.06):
    sh = s.shapes.add_shape(MSO.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if lc: sh.line.color.rgb = lc; sh.line.width = lw
    else:  sh.line.fill.background()
    sh.adjustments[0] = adj
    return sh

def oval(s, x, y, w, h, fc):
    sh = s.shapes.add_shape(MSO.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.fill.background()
    return sh

def line_h(s, x, y, w, c, lw=Pt(1)):
    sh = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    sh.line.color.rgb = c; sh.line.width = lw; return sh

def line_v(s, x, y, h, c, lw=Pt(1)):
    sh = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x), Inches(y+h))
    sh.line.color.rgb = c; sh.line.width = lw; return sh

def txt(s, text, x, y, w, h, size=12, c=None, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    if c is None: c = DK_TXT
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = c
    r.font.bold = bold; r.font.name = 'Calibri'
    return tb

def txt_ml(s, lines, x, y, w, h, size=12, c=None, bold=False, align=PP_ALIGN.LEFT):
    """Multi-line textbox (list of strings)"""
    if c is None: c = DK_TXT
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = c
        r.font.bold = bold; r.font.name = 'Calibri'
    return tb

def card(s, x, y, w, h):
    rrect(s, x, y, w, h, WHITE, LT_GRAY, Pt(0.75))

def header(s, num, label):
    rect(s, 0, 0, 10, 1.05, NAVY)
    txt(s, f"{num}  {label}", 0.5, 0.18, 9, 0.7, size=25, c=WHITE, bold=True)

def badge(s, x, y, w, h, label, fc=None):
    if fc is None: fc = NAVY
    rrect(s, x, y, w, h, fc, adj=0.12)
    txt(s, label, x, y, w, h, size=11, c=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# 슬라이드 1: 타이틀
# ════════════════════════════════════════════
s = slide(); bg(s, NAVY)
rect(s, 0, 3.5, 10, 2.125, BLUE)
txt(s, "한성과학고등학교", 0.5, 0.85, 9, 0.55, size=17, c=LT_BLUE, align=PP_ALIGN.CENTER)
txt(s, "입시 결과 분석 대시보드", 0.5, 1.45, 9, 1.0, size=42, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "연수 가이드", 0.5, 2.55, 9, 0.55, size=22, c=GOLD, align=PP_ALIGN.CENTER)
txt(s, "대학교별 내신 등급 분포 시각화  ·  학생 상담 실전 활용법", 0.5, 3.7, 9, 0.5,
    size=13, c=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 2: 목차
# ════════════════════════════════════════════
s = slide(); bg(s, LT_BG)
txt(s, "연수 구성", 0.5, 0.28, 9, 0.65, size=30, c=NAVY, bold=True)
rect(s, 0.5, 0.95, 9, 0.03, LT_BLUE)

items = [
    ("01", "기본 필터 조작",      "연도 · 전형 · 계열 · 성적 기준 선택"),
    ("02", "박스 플롯 읽는 법",   "분포 구조 · 수염 · 개별 점의 의미"),
    ("03", "범례 구분 이해",      "합격 상태 · 추가합격 · 미응시 구분"),
    ("04", "드릴다운 모달 활용",  "연도별 · 학과별 · 전형별 상세 분석"),
    ("05", "내 등급 기준선",      "학생 상담 실전 활용법"),
    ("06", "해석 시 주의사항",    "데이터 신뢰도 판단 기준"),
    ("07", "기능 전체 체크리스트", "구현된 기능 목록 한눈에 확인"),
]
for i, (num, title, desc) in enumerate(items):
    y = 1.08 + i * 0.62
    oval(s, 0.5, y + 0.06, 0.42, 0.42, NAVY)
    txt(s, num,   0.5,  y+0.06, 0.42, 0.42, size=11, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, 1.08, y+0.04, 3.0,  0.32, size=13, c=DK_TXT, bold=True)
    txt(s, desc,  4.2,  y+0.04, 5.5,  0.32, size=12, c=GRAY)

# ════════════════════════════════════════════
# 슬라이드 3: 기본 필터 조작
# ════════════════════════════════════════════
s = slide(); bg(s, LT_BG); header(s, "01", "기본 필터 조작")

filters = [
    ("연도 선택",      "다중 선택 가능\n여러 연도 합산으로 경향 비교"),
    ("학년 필터",      "졸업 기준 학년 구분\n(1 ~ 3학년)"),
    ("전형 탭",        "학생부종합 / 논술\n특기자전형 / 특별전형"),
    ("계열 전환",      "일반계열\n메디컬계열"),
    ("성적 기준",      "국영수과 등평\n↕  전교과 등평"),
    ("합격 없는 대학", "표시 여부를\n토글로 제어 가능"),
]
xs = [0.3, 3.5, 6.7];  ys = [1.18, 3.05]
for i, (name, desc) in enumerate(filters):
    x = xs[i % 3];  y = ys[i // 3]
    card(s, x, y, 3.0, 1.68)
    txt(s, name, x+0.18, y+0.12, 2.65, 0.36, size=13.5, c=NAVY, bold=True)
    rect(s, x+0.18, y+0.52, 2.65, 0.025, LT_BLUE)
    txt(s, desc, x+0.18, y+0.6, 2.65, 0.95, size=11.5, c=GRAY, wrap=True)

# ════════════════════════════════════════════
# 슬라이드 4: 박스 플롯 읽는 법
# ════════════════════════════════════════════
s = slide(); bg(s, LT_BG); header(s, "02", "박스 플롯 읽는 법")

# 다이어그램 좌표
cx = 2.0; bw = 0.7
wT = 1.28; q3 = 1.88; med = 2.52; q1 = 3.2; wB = 3.82

# 수염 세로선
line_v(s, cx, wT,  q3-wT,   BLUE, Pt(2.5))
line_v(s, cx, q1,  wB-q1,   BLUE, Pt(2.5))
# 수염 가로 캡
line_h(s, cx-bw, wT, bw*2, BLUE, Pt(2.5))
line_h(s, cx-bw, wB, bw*2, BLUE, Pt(2.5))
# 박스
rrect(s, cx-bw, q3, bw*2, q1-q3, LT_BLUE, BLUE, Pt(2.5), adj=0.04)
# 중앙값 선
line_h(s, cx-bw, med, bw*2, NAVY, Pt(3.5))
# 개별 점
for px, py in [(cx-.18,1.55),(cx+.12,2.05),(cx-.08,2.22),
               (cx+.16,2.75),(cx-.05,3.0),(cx+.1,3.42),(cx-.15,3.58)]:
    oval(s, px-.07, py-.07, 0.14, 0.14, BLUE)

# 레이블
lx = cx + bw + 0.15
for y, title, sub, hl in [
    (wT-.22, "위 수염 끝",     "Q3 + 1.5×IQR 이하 실제 최댓값",  False),
    (q3-.14, "Q3  (75% 위치)", "중간 50% 구간 상단 경계",          False),
    (med-.14,"중앙값 (Median)","50% 위치 · 박스 안 가로선",         True ),
    (q1-.14, "Q1  (25% 위치)", "중간 50% 구간 하단 경계",          False),
    (wB-.06, "아래 수염 끝",   "Q1 − 1.5×IQR 이상 실제 최솟값",  False),
]:
    txt(s, title, lx, y,     2.65, 0.28, size=12, c=NAVY if hl else DK_TXT, bold=True)
    txt(s, sub,   lx, y+0.3, 2.65, 0.24, size=10.5, c=GRAY)

# IQR 레이블
txt(s, "IQR\n(박스 높이)", cx-bw-1.1, (q1+q3)/2-.4, 0.88, 0.5,
    size=11, c=NAVY, bold=True, align=PP_ALIGN.CENTER)

# 오른쪽 요점 카드
card(s, 5.0, 1.18, 4.65, 4.1)
txt(s, "핵심 포인트", 5.2, 1.28, 4.3, 0.4, size=14, c=NAVY, bold=True)
pts = [
    ("박스 높이 = IQR",  "중간 50% 학생의 등급 분포 범위"),
    ("수염 범위",        "IQR × 1.5 내 실제 데이터까지"),
    ("개별 점 (●)",      "실제 지원 학생 한 명 한 명"),
    ("Y축 방향 주의",    "숫자가 낮을수록 우수한 성적"),
    ("X축 괄호 표기",    "(최종합격수 / 총지원수)"),
]
for i, (t, d) in enumerate(pts):
    py = 1.85 + i*0.65
    oval(s, 5.2, py+0.05, 0.2, 0.2, NAVY)
    txt(s, t, 5.55, py,      3.9, 0.28, size=12.5, c=DK_TXT, bold=True)
    txt(s, d, 5.55, py+0.3,  3.9, 0.26, size=11,   c=GRAY)

# ════════════════════════════════════════════
# 슬라이드 5: 범례 구분 이해
# ════════════════════════════════════════════
s = slide(); bg(s, LT_BG); header(s, "03", "범례 구분 이해")

legends = [
    ("2196F3", "최종합격",             "최종 합격한 학생의 성적 분포"),
    ("E67E22", "└ 추가합격  (PC 전용)", "추가합격 토글 ON 시 주황색 구분 · 최초합격보다 등급이 낮을 수 있음"),
    ("C0392B", "1차합격_최종탈락",      "1차 합격 후 최종 불합격 (등록 포기, 타교 합격 등)"),
    ("AAAAAA", "└ ✕  미응시",           "2차 시험 미응시 — 박스 통계에서 제외, 해석 시 참고만 권장"),
    ("BBBBBB", "1차탈락  (참고)",        "1차 탈락 학생 데이터 — 지원 가능 하한선 파악 시 참고용"),
]
for i, (hx, label, desc) in enumerate(legends):
    y = 1.18 + i*0.85
    card(s, 0.4, y, 9.2, 0.73)
    oval(s, 0.65, y+0.19, 0.34, 0.34, rgb(hx))
    txt(s, label, 1.15, y+0.09, 3.2, 0.3,  size=13,   c=DK_TXT, bold=True)
    txt(s, desc,  1.15, y+0.39, 8.2, 0.28, size=11.5, c=GRAY)

# ════════════════════════════════════════════
# 슬라이드 6: 드릴다운 모달
# ════════════════════════════════════════════
s = slide(); bg(s, LT_BG); header(s, "04", "드릴다운 모달 활용")

rrect(s, 0.4, 1.15, 9.2, 0.6, LT_BLUE, BLUE, Pt(0.75))
txt(s, "박스 플롯에서 대학명을 클릭하면 상세 분석 팝업이 열립니다",
    0.6, 1.27, 8.8, 0.38, size=13, c=NAVY, bold=True)

tabs = [
    ("탭 1", "연도별",  ["연도별 합격선 변화 추이 확인", "특정 연도 이상치 여부 파악", "최근 3년 트렌드 중심 해석"]),
    ("탭 2", "학과별",  ["모집단위별 등급 분포 비교", "유리한 학과 탐색", "학과 간 합격선 격차 확인"]),
    ("탭 3", "전형별",  ["전형명칭별 합격선 비교", "학생부종합 vs 논술 등\n같은 대학 내 전형 간 전략 비교"]),
]
for i, (num, title, pts) in enumerate(tabs):
    x = 0.4 + i*3.15
    card(s, x, 1.9, 3.0, 3.45)
    rrect(s, x+0.1, 2.0, 2.8, 0.5, NAVY, adj=0.1)
    txt(s, f"{num}  {title}", x+0.1, 2.0, 2.8, 0.5, size=14, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, pt in enumerate(pts):
        txt_ml(s, ["• " + ln for ln in pt.split("\n")],
               x+0.2, 2.65+j*0.85, 2.6, 0.75, size=12, c=DK_TXT)

# ════════════════════════════════════════════
# 슬라이드 7: 내 등급 기준선
# ════════════════════════════════════════════
s = slide(); bg(s, LT_BG); header(s, "05", "내 등급 기준선 활용")

card(s, 0.4, 1.18, 4.4, 4.1)
txt(s, "사용 방법", 0.6, 1.28, 4.0, 0.4, size=14, c=NAVY, bold=True)

steps = [
    ("① 등급 입력",      "상단 입력칸에 학생의 등평 입력"),
    ("② 빨간 점선 확인", "차트에 빨간 수평 기준선이 표시됨"),
    ("③ 요약 확인",      "하단에 안정권·적정권·도전권\n대학 목록이 자동 생성"),
]
for i, (badge_txt, desc) in enumerate(steps):
    sy = 1.88 + i*1.1
    rrect(s, 0.6, sy, 1.55, 0.36, NAVY, adj=0.1)
    txt(s, badge_txt, 0.6, sy, 1.55, 0.36, size=10.5, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt_ml(s, desc.split("\n"), 0.6, sy+0.44, 4.0, 0.55, size=11.5, c=GRAY)

card(s, 5.05, 1.18, 4.55, 4.1)
txt(s, "상담 활용 팁", 5.25, 1.28, 4.15, 0.4, size=14, c=NAVY, bold=True)

tips = [
    "등급 입력 → 합격 가능성 시각적 파악",
    "최근 2~3년도만 선택해 최신 트렌드 반영",
    "학과별 탭으로 같은 대학 내 유리한 학과 발견",
    "전형별 탭으로 학생부종합 vs 논술 전략 비교",
    "1차탈락 데이터 참고해 지원 가능 하한선 설정",
]
for i, tip in enumerate(tips):
    txt(s, "• " + tip, 5.25, 1.88+i*0.65, 4.15, 0.55, size=11.5, c=DK_TXT)

# ════════════════════════════════════════════
# 슬라이드 8: 해석 시 주의사항
# ════════════════════════════════════════════
s = slide(); bg(s, LT_BG); header(s, "06", "해석 시 주의사항")

warnings = [
    ("데이터 수가 적은 대학",  "샘플이 적으면 통계적 신뢰도 낮음 — 점의 수를 확인하고 해석"),
    ("미응시 X 마커",         "2차 시험 미응시자는 실질 경쟁 데이터가 아니므로 합격선 산출 시 제외"),
    ("KAIST 분리 표기",       "KAIST(일반전형)와 KAIST(창의도전)은 별개 전형으로 분리 표기됨"),
    ("추가합격 해석 주의",     "추가합격자 등급이 최초합격보다 낮을 수 있어 단순 비교 금물"),
    ("연도별 전형 변화",       "해마다 선발 방식이 달라지므로 최신 연도 데이터를 우선 참고"),
]
for i, (title, desc) in enumerate(warnings):
    y = 1.2 + i*0.85
    card(s, 0.4, y, 9.2, 0.73)
    oval(s, 0.62, y+0.19, 0.34, 0.34, GOLD)
    txt(s, "!",  0.62, y+0.19, 0.34, 0.34, size=13, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, 1.12, y+0.09, 2.9, 0.3,  size=12.5, c=DK_TXT, bold=True)
    txt(s, desc,  1.12, y+0.4,  8.3, 0.27, size=11,   c=GRAY)

# ════════════════════════════════════════════
# 슬라이드 9: 기능 전체 체크리스트
# ════════════════════════════════════════════
s = slide(); bg(s, LT_BG); header(s, "07", "기능 전체 체크리스트")

# 왼쪽 열
left_items = [
    ("필터 · 조건 설정",      True, [
        "연도 다중 선택 / 학년 필터",
        "전형 탭 (학생부종합·논술·특기자·특별)",
        "계열 전환 (일반 ↔ 메디컬)",
        "성적 기준 전환 (국영수과 ↔ 전교과)",
        "합격 없는 대학 표시 토글",
    ]),
    ("추가합격 구분  (PC 전용)", True, [
        "토글 ON 시 추가합격 점 주황색 표시",
        "최초합격 / 추가합격 범례 구분",
    ]),
    ("미응시 별도 표시", True, [
        "✕ 마커로 박스와 분리 표시",
        "박스 통계(중앙값·IQR)에서 제외",
    ]),
    ("내 등급 기준선", True, [
        "빨간 점선 수평선 표시",
        "안정권 · 적정권 · 도전권 자동 분류",
    ]),
]

right_items = [
    ("드릴다운 모달", True, [
        "대학명 클릭 → 연도별 · 학과별 · 전형별 탭",
        "학과 키워드 검색 (다중 OR)",
        "줌 · 팬 (휠 스크롤 + 드래그)",
    ]),
    ("마우스 호버 정보", True, [
        "학과명 · 합격 구분 (최초 / 추가)",
        "⚠ 수능최저 미충족 표시",
        "예비번호 표시",
        "추가합격 상세 정보",
        "등록 여부",
    ]),
    ("접근 제어", True, [
        "Google 로그인 인증",
        "허가된 교사 계정만 접근 가능",
    ]),
    ("모바일 지원", True, [
        "핀치줌 / 드래그 스크롤",
        "반응형 레이아웃",
    ]),
]

def check_section(s, x, y, group_label, done, subitems):
    check = "✅" if done else "☐"
    txt(s, f"{check}  {group_label}", x, y, 4.3, 0.3, size=12.5, c=NAVY, bold=True)
    for i, item in enumerate(subitems):
        txt(s, f"   · {item}", x, y + 0.3 + i*0.25, 4.3, 0.25, size=10.5, c=GRAY)
    return y + 0.3 + len(subitems)*0.25 + 0.15

# 구분선
line_v(s, 5.1, 1.12, 4.35, LT_GRAY, Pt(1))

ly = 1.15
for group_label, done, subitems in left_items:
    ly = check_section(s, 0.45, ly, group_label, done, subitems)

ry = 1.15
for group_label, done, subitems in right_items:
    ry = check_section(s, 5.25, ry, group_label, done, subitems)

# ════════════════════════════════════════════
# 슬라이드 10: 마무리
# ════════════════════════════════════════════
s = slide(); bg(s, NAVY)
rect(s, 0, 2.18, 10, 1.27, BLUE)
txt(s, "감사합니다", 0.5, 0.9, 9, 1.05, size=48, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "한성과학고등학교 입시 결과 분석 대시보드", 0.5, 2.35, 9, 0.5,
    size=17, c=WHITE, align=PP_ALIGN.CENTER)
txt(s, "궁금한 기능이나 추가 요청이 있으시면 언제든 말씀해 주세요",
    0.5, 3.65, 9, 0.5, size=14, c=LT_BLUE, align=PP_ALIGN.CENTER)

# ── 저장 ───────────────────────────────────
prs.save(OUT)
print(f"저장 완료: {OUT}")
